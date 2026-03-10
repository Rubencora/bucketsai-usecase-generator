"""
BucketsAI PPTX Template Generator
==================================
Generates a 12-slide professional PPTX template for pptx-automizer.
Each text element is named via shape.name and contains {{placeholder}} text.

Brand palette (BucketsAI):
  BRAND_BLUE   #4470DC  — CTAs, highlights, accent bars
  BLUE_MEDIUM  #6B8FE8  — secondary elements
  BLUE_LIGHT   #E4EBF8  — card backgrounds
  PAGE_BG      #EEF0F6  — page background
  NAVY_MID     #2E3E5C  — dark panels
  NAVY_LIGHT   #3A4D6E  — cards on dark bg
  DARK_NAVY    #0C1628  — TEXT ONLY
  TEXT_MUTED   #5A6A88  — body text
  GRAY_MID     #8A96B0  — captions
  BORDER       #D8DCE8  — borders
  GREEN        #22B573  — positive
  RED          #E74C3C  — negative
  ORANGE       #FF6B35  — accent only
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from lxml import etree


# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------

BRAND_BLUE = RGBColor(0x44, 0x70, 0xDC)
BLUE_MEDIUM = RGBColor(0x6B, 0x8F, 0xE8)
BLUE_LIGHT = RGBColor(0xE4, 0xEB, 0xF8)
PAGE_BG = RGBColor(0xEE, 0xF0, 0xF6)
NAVY_MID = RGBColor(0x2E, 0x3E, 0x5C)
NAVY_LIGHT = RGBColor(0x3A, 0x4D, 0x6E)
DARK_NAVY = RGBColor(0x0C, 0x16, 0x28)
TEXT_MUTED = RGBColor(0x5A, 0x6A, 0x88)
GRAY_MID = RGBColor(0x8A, 0x96, 0xB0)
BORDER = RGBColor(0xD8, 0xDC, 0xE8)
GREEN = RGBColor(0x22, 0xB5, 0x73)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_LIGHT = RGBColor(0xFF, 0xF5, 0xF5)
GREEN_LIGHT = RGBColor(0xF0, 0xFB, 0xF6)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

FONT = "Arial"

OUTPUT_PATH = Path(
    "/Users/rubencordoba/Library/CloudStorage/"
    "GoogleDrive-ruben@hoytrabajas.com/My Drive/"
    "Buckets AI - Hoytrabajas - Talentropy. 2026/"
    "5. AI Projects/UseCase BucketsAI/"
    "bucketsai-usecase-generator/templates/deck_template.pptx"
)


# ---------------------------------------------------------------------------
# Helpers — geometry, shadows, text
# ---------------------------------------------------------------------------


def _hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def add_shadow(shape) -> None:
    """Attach a subtle outer drop shadow to a shape via XML."""
    spPr = shape._element.spPr
    # Remove existing effectLst if present
    for existing in spPr.findall(qn("a:effectLst")):
        spPr.remove(existing)

    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
    outerShdw.set("blurRad", "50800")   # 4pt
    outerShdw.set("dist", "25400")      # 2pt
    outerShdw.set("dir", "5400000")     # 90°
    outerShdw.set("algn", "tl")
    outerShdw.set("rotWithShape", "0")
    srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"))
    srgbClr.set("val", "000000")
    alpha = etree.SubElement(srgbClr, qn("a:alpha"))
    alpha.set("val", "15000")           # 15 %


def set_no_line(shape) -> None:
    """Remove border/line from shape."""
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    noFill = ln.find(qn("a:noFill"))
    if noFill is None:
        etree.SubElement(ln, qn("a:noFill"))


def set_line(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    """Set a solid border on a shape."""
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(width_pt * 12700)))
    # Clear existing fill
    for child in list(ln):
        ln.remove(child)
    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", _hex(color))


def set_rounded_corners(shape, adj_val: int = 25000) -> None:
    """Enforce roundRect preset geometry with an adjustable corner radius."""
    sp = shape._element
    spPr = sp.spPr
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "roundRect")
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        gd = avLst.find(qn("a:gd"))
        if gd is None:
            gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {adj_val}")


def fill_shape(shape, color: RGBColor) -> None:
    """Set a solid fill colour on a shape."""
    spPr = shape._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is not None:
        spPr.remove(solidFill)
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", _hex(color))


def add_textbox(
    slide,
    name: str,
    placeholder: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 12,
    bold: bool = False,
    color: RGBColor = DARK_NAVY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> object:
    """Add a named textbox with placeholder text."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.name = name
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = placeholder
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rounded_rect(
    slide,
    name: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor = WHITE,
    shadow: bool = True,
    adj: int = 25000,
    line_color: Optional[RGBColor] = None,
) -> object:
    """Add a rounded-rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.name = name
    fill_shape(shape, fill)
    set_rounded_corners(shape, adj)
    if line_color:
        set_line(shape, line_color, 0.75)
    else:
        set_no_line(shape)
    if shadow:
        add_shadow(shape)
    return shape


def add_rect(
    slide,
    name: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor = BRAND_BLUE,
) -> object:
    """Add a plain rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches as I

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.name = name
    fill_shape(shape, fill)
    set_no_line(shape)
    return shape


def add_circle(
    slide,
    name: str,
    left: float,
    top: float,
    diameter: float,
    fill: RGBColor = BRAND_BLUE,
    text: str = "",
    font_size: int = 10,
    text_color: RGBColor = WHITE,
) -> object:
    """Add a circle (oval with equal w/h) shape with optional centred text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter),
    )
    shape.name = name
    fill_shape(shape, fill)
    set_no_line(shape)
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = text_color
    return shape


def set_slide_background(slide, color: RGBColor) -> None:
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_label_in_shape(
    shape,
    text: str,
    font_size: int = 11,
    bold: bool = False,
    color: RGBColor = DARK_NAVY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Write text into an existing shape's text frame (replaces content)."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_footer(slide, url_text: str = "buckets-ai.com") -> None:
    """Add a subtle footer line with URL."""
    footer = add_textbox(
        slide, "footer_url", url_text,
        left=0.5, top=7.1, width=4.0, height=0.3,
        font_size=9, color=GRAY_MID,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def build_slide_01_cover(prs: Presentation) -> None:
    """Slide 1: Cover — dark navy background, hero text, decorative circles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, NAVY_MID)

    # Decorative semi-transparent circles (emulate with low-opacity shapes)
    for name, l, t, d in [
        ("deco_circle_1", 10.5, -0.8, 3.5),
        ("deco_circle_2", 11.8, 0.8, 2.0),
        ("deco_circle_3", -0.5, 5.5, 2.5),
    ]:
        c = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(l), Inches(t), Inches(d), Inches(d),
        )
        c.name = name
        spPr = c._element.spPr
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", _hex(BRAND_BLUE))
        alpha = etree.SubElement(srgbClr, qn("a:alpha"))
        alpha.set("val", "8000")  # 8 %
        set_no_line(c)

    # Thin blue accent line on left
    accent = add_rect(slide, "cover_accent_bar", 0.5, 1.8, 0.06, 3.5, fill=BRAND_BLUE)

    # Company label
    add_textbox(
        slide, "cover_empresa", "{{empresa}}",
        left=0.7, top=1.8, width=6.0, height=0.4,
        font_size=13, color=BLUE_MEDIUM, bold=False,
    )

    # Main tagline
    add_textbox(
        slide, "cover_tagline", "{{tagline}}",
        left=0.7, top=2.3, width=9.0, height=2.0,
        font_size=36, bold=True, color=WHITE,
    )

    # URL
    add_textbox(
        slide, "cover_url", "buckets-ai.com",
        left=0.7, top=6.8, width=3.0, height=0.35,
        font_size=10, color=GRAY_MID,
    )

    # Logo placeholder (image area marked with a rounded rect)
    logo_ph = add_rounded_rect(
        slide, "cover_logo",
        left=11.0, top=6.5, width=1.8, height=0.7,
        fill=NAVY_LIGHT, shadow=False,
    )
    add_label_in_shape(logo_ph, "LOGO", font_size=9, color=GRAY_MID, align=PP_ALIGN.CENTER)


def build_slide_02_problem(prs: Presentation) -> None:
    """Slide 2: Problem — two-column card layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Headline
    add_textbox(
        slide, "problem_headline", "{{problem_headline}}",
        left=0.5, top=0.35, width=12.33, height=0.55,
        font_size=24, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER,
    )

    # ---- Left card (blue light) ----
    card_l = add_rounded_rect(
        slide, "problem_left_card",
        left=0.5, top=1.1, width=5.9, height=5.0,
        fill=BLUE_LIGHT, shadow=True,
    )
    # Blue top accent bar
    accent_l = add_rounded_rect(
        slide, "problem_left_accent",
        left=0.5, top=1.1, width=5.9, height=0.12,
        fill=BRAND_BLUE, shadow=False, adj=10000,
    )
    add_textbox(
        slide, "problem_left_title", "Tus equipos tienen:",
        left=0.7, top=1.3, width=5.5, height=0.4,
        font_size=14, bold=True, color=DARK_NAVY,
    )
    for i in range(1, 5):
        add_textbox(
            slide, f"problem_has_{i}", f"{{{{problem_has_{i}}}}}",
            left=0.8, top=1.75 + (i - 1) * 0.9, width=5.3, height=0.75,
            font_size=12, color=TEXT_MUTED,
        )

    # ---- Right card (light red) ----
    card_r = add_rounded_rect(
        slide, "problem_right_card",
        left=6.93, top=1.1, width=5.9, height=5.0,
        fill=RED_LIGHT, shadow=True,
    )
    accent_r = add_rounded_rect(
        slide, "problem_right_accent",
        left=6.93, top=1.1, width=5.9, height=0.12,
        fill=RED, shadow=False, adj=10000,
    )
    add_textbox(
        slide, "problem_right_title", "Pero en la ejecucion real:",
        left=7.1, top=1.3, width=5.5, height=0.4,
        font_size=14, bold=True, color=DARK_NAVY,
    )
    for i in range(1, 5):
        add_textbox(
            slide, f"problem_but_{i}", f"{{{{problem_but_{i}}}}}",
            left=7.2, top=1.75 + (i - 1) * 0.9, width=5.3, height=0.75,
            font_size=12, color=TEXT_MUTED,
        )

    # Bottom bar
    bar = add_rect(slide, "problem_closing_bar", 0.5, 6.3, 12.33, 0.6, fill=NAVY_MID)
    add_textbox(
        slide, "problem_closing", "{{problem_closing}}",
        left=0.7, top=6.32, width=11.9, height=0.5,
        font_size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER,
    )

    add_footer(slide)


def build_slide_03_depth(prs: Presentation) -> None:
    """Slide 3: Depth — split layout, 5 reason cards on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PAGE_BG)

    # Left half — navy panel
    left_panel = add_rect(slide, "depth_left_panel", 0.0, 0.0, 5.5, 7.5, fill=NAVY_MID)

    add_textbox(
        slide, "depth_headline", "{{depth_headline}}",
        left=0.5, top=2.8, width=4.5, height=2.0,
        font_size=30, bold=True, color=WHITE,
    )

    # Decorative accent line
    add_rect(slide, "depth_accent", 0.5, 2.55, 0.08, 0.22, fill=BRAND_BLUE)

    # Right half — 5 stacked reason cards
    card_h = 1.15
    card_top_start = 0.35
    for i in range(1, 6):
        top = card_top_start + (i - 1) * (card_h + 0.08)
        card = add_rounded_rect(
            slide, f"depth_reason_card_{i}",
            left=5.8, top=top, width=7.03, height=card_h,
            fill=WHITE, shadow=True,
        )
        # Blue left accent bar
        add_rect(
            slide, f"depth_reason_accent_{i}",
            left=5.8, top=top, width=0.07, height=card_h,
            fill=BRAND_BLUE,
        )
        # Number circle
        add_circle(
            slide, f"depth_reason_num_{i}",
            left=6.0, top=top + 0.28, diameter=0.55,
            fill=BRAND_BLUE, text=str(i), font_size=11,
        )
        # Text
        add_textbox(
            slide, f"depth_reason_{i}", f"{{{{depth_reason_{i}}}}}",
            left=6.7, top=top + 0.22, width=6.0, height=0.75,
            font_size=12, color=DARK_NAVY,
        )

    add_footer(slide)


def build_slide_04_solution(prs: Presentation) -> None:
    """Slide 4: Solution — left text + right navy capability panel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PAGE_BG)

    # Left side content
    add_textbox(
        slide, "solution_headline", "{{solution_headline}}",
        left=0.5, top=1.2, width=6.0, height=1.4,
        font_size=28, bold=True, color=DARK_NAVY,
    )
    add_textbox(
        slide, "solution_description", "{{solution_description}}",
        left=0.5, top=2.8, width=5.8, height=2.5,
        font_size=14, color=TEXT_MUTED,
    )

    # Blue accent underline for headline
    add_rect(slide, "solution_headline_accent", 0.5, 2.65, 2.5, 0.05, fill=BRAND_BLUE)

    # Right navy panel
    right_panel = add_rounded_rect(
        slide, "solution_right_panel",
        left=7.0, top=0.6, width=5.83, height=6.3,
        fill=NAVY_MID, shadow=True,
    )
    add_textbox(
        slide, "solution_panel_title", "Le permite a los equipos:",
        left=7.3, top=0.9, width=5.2, height=0.5,
        font_size=14, bold=True, color=WHITE,
    )
    # Blue separator under panel title
    add_rect(slide, "solution_panel_separator", 7.3, 1.45, 5.2, 0.04, fill=BLUE_MEDIUM)

    for i in range(1, 6):
        top = 1.6 + (i - 1) * 0.95
        # Check mark circle
        add_circle(
            slide, f"solution_check_{i}",
            left=7.3, top=top + 0.08, diameter=0.4,
            fill=BRAND_BLUE, text="✓", font_size=10,
        )
        add_textbox(
            slide, f"solution_cap_{i}", f"{{{{solution_cap_{i}}}}}",
            left=7.85, top=top, width=4.8, height=0.8,
            font_size=12, color=WHITE,
        )

    add_footer(slide)


def build_slide_05_how_it_works(prs: Presentation) -> None:
    """Slide 5: How it works — 3 step cards + bottom strip."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, "how_title", "Como funciona",
        left=0.5, top=0.3, width=12.33, height=0.55,
        font_size=26, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, "how_subtitle", "Listo para operar en horas",
        left=0.5, top=0.85, width=12.33, height=0.4,
        font_size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )

    # Blue accent underline centred under title
    add_rect(slide, "how_title_accent", 5.8, 0.82, 1.73, 0.05, fill=BRAND_BLUE)

    step_configs = [
        ("how_step_card_1", BRAND_BLUE, "01", "how_step_1"),
        ("how_step_card_2", GREEN, "02", "how_step_2"),
        ("how_step_card_3", ORANGE, "03", "how_step_3"),
    ]
    card_w = 3.8
    card_gap = 0.26
    card_left_start = 0.5

    for idx, (card_name, accent_color, num, text_name) in enumerate(step_configs):
        left = card_left_start + idx * (card_w + card_gap)
        # Card body
        card = add_rounded_rect(
            slide, card_name,
            left=left, top=1.5, width=card_w, height=4.6,
            fill=WHITE, shadow=True,
        )
        # Coloured top accent
        add_rounded_rect(
            slide, f"{card_name}_accent",
            left=left, top=1.5, width=card_w, height=0.15,
            fill=accent_color, shadow=False, adj=10000,
        )
        # Number circle
        add_circle(
            slide, f"{card_name}_num",
            left=left + 0.2, top=1.75, diameter=0.65,
            fill=accent_color, text=num, font_size=13,
        )
        # Step text (placeholder)
        add_textbox(
            slide, text_name, f"{{{{{text_name}}}}}",
            left=left + 0.2, top=2.55, width=3.4, height=3.3,
            font_size=12, color=TEXT_MUTED,
        )

    # Bottom strip
    strip = add_rect(slide, "how_strip", 0.0, 6.25, 13.33, 1.0, fill=BLUE_LIGHT)
    add_textbox(
        slide, "how_not_items", "{{how_not_items}}",
        left=0.5, top=6.35, width=12.33, height=0.75,
        font_size=11, color=DARK_NAVY, align=PP_ALIGN.CENTER,
    )

    add_footer(slide)


def build_slide_06_demo(prs: Presentation) -> None:
    """Slide 6: Demo — headline left + chat mockup card right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PAGE_BG)

    add_textbox(
        slide, "demo_headline", "{{demo_headline}}",
        left=0.5, top=1.2, width=5.5, height=1.4,
        font_size=26, bold=True, color=DARK_NAVY,
    )
    add_textbox(
        slide, "demo_description", "{{demo_description}}",
        left=0.5, top=2.8, width=5.5, height=2.5,
        font_size=14, color=TEXT_MUTED,
    )
    add_rect(slide, "demo_headline_accent", 0.5, 2.65, 2.0, 0.05, fill=BRAND_BLUE)

    # Chat mockup outer card
    chat_card = add_rounded_rect(
        slide, "demo_chat_card",
        left=6.5, top=0.8, width=6.33, height=6.3,
        fill=WHITE, shadow=True,
    )
    # Navy header bar
    header = add_rounded_rect(
        slide, "demo_chat_header",
        left=6.5, top=0.8, width=6.33, height=0.65,
        fill=NAVY_MID, shadow=False, adj=8000,
    )
    add_textbox(
        slide, "demo_chat_header_label", "BucketsAI Chat",
        left=6.7, top=0.88, width=4.0, height=0.4,
        font_size=11, bold=True, color=WHITE,
    )
    # Status dot
    add_circle(
        slide, "demo_status_dot",
        left=12.38, top=0.98, diameter=0.2,
        fill=GREEN,
    )

    # User message bubble (blue light)
    user_bubble = add_rounded_rect(
        slide, "demo_user_bubble",
        left=8.5, top=1.7, width=4.0, height=1.1,
        fill=BLUE_LIGHT, shadow=False,
    )
    add_textbox(
        slide, "demo_user_message", "{{demo_user_message}}",
        left=8.6, top=1.78, width=3.8, height=0.9,
        font_size=11, color=DARK_NAVY,
    )

    # AI response bubble (white with border)
    ai_bubble = add_rounded_rect(
        slide, "demo_ai_bubble",
        left=6.7, top=3.1, width=5.5, height=2.2,
        fill=WHITE, shadow=False,
        line_color=BORDER,
    )
    add_textbox(
        slide, "demo_ai_response", "{{demo_ai_response}}",
        left=6.85, top=3.22, width=5.2, height=1.9,
        font_size=11, color=TEXT_MUTED,
    )

    add_footer(slide)


def build_slide_07_dimensions(prs: Presentation) -> None:
    """Slide 7: Two value dimensions — side-by-side cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, "dimensions_title", "Dos dimensiones de valor",
        left=0.5, top=0.3, width=12.33, height=0.55,
        font_size=26, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, "dimensions_accent", 5.9, 0.82, 1.53, 0.05, fill=BRAND_BLUE)

    card_configs = [
        ("dim_a", BRAND_BLUE, 0.5),
        ("dim_b", ORANGE, 6.9),
    ]
    for prefix, accent_color, left in card_configs:
        card = add_rounded_rect(
            slide, f"{prefix}_card",
            left=left, top=1.1, width=5.9, height=5.8,
            fill=BLUE_LIGHT if accent_color == BRAND_BLUE else RGBColor(0xFF, 0xF8, 0xF4),
            shadow=True,
        )
        # Accent top bar
        add_rounded_rect(
            slide, f"{prefix}_accent_bar",
            left=left, top=1.1, width=5.9, height=0.15,
            fill=accent_color, shadow=False, adj=10000,
        )
        # Label badge
        badge = add_rounded_rect(
            slide, f"{prefix}_badge",
            left=left + 0.25, top=1.4, width=1.3, height=0.35,
            fill=accent_color, shadow=False, adj=50000,
        )
        add_textbox(
            slide, f"{prefix}_label", f"{{{{{prefix}_label}}}}",
            left=left + 0.27, top=1.42, width=1.26, height=0.3,
            font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, f"{prefix}_titulo", f"{{{{{prefix}_titulo}}}}",
            left=left + 0.25, top=1.9, width=5.4, height=0.75,
            font_size=18, bold=True, color=DARK_NAVY,
        )
        add_textbox(
            slide, f"{prefix}_rol", f"{{{{{prefix}_rol}}}}",
            left=left + 0.25, top=2.7, width=5.4, height=0.4,
            font_size=11, color=accent_color, bold=True,
        )
        add_textbox(
            slide, f"{prefix}_desc", f"{{{{{prefix}_desc}}}}",
            left=left + 0.25, top=3.2, width=5.4, height=3.4,
            font_size=12, color=TEXT_MUTED,
        )

    add_footer(slide)


def build_slide_08_before_after(prs: Presentation) -> None:
    """Slide 8: Before/After comparison columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, "beforeafter_title", "{{beforeafter_title}}",
        left=0.5, top=0.3, width=12.33, height=0.55,
        font_size=24, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER,
    )

    col_configs = [
        ("before", RED, RED_LIGHT, 0.5, "ANTES", "✕"),
        ("after", GREEN, GREEN_LIGHT, 6.93, "DESPUES", "✓"),
    ]
    for prefix, accent, bg, left, label_text, icon in col_configs:
        # Column background
        col_card = add_rounded_rect(
            slide, f"{prefix}_column",
            left=left, top=1.0, width=5.9, height=5.9,
            fill=bg, shadow=True,
        )
        add_rounded_rect(
            slide, f"{prefix}_accent_top",
            left=left, top=1.0, width=5.9, height=0.15,
            fill=accent, shadow=False, adj=10000,
        )
        # Label
        add_textbox(
            slide, f"{prefix}_label", label_text,
            left=left + 0.25, top=1.25, width=2.0, height=0.45,
            font_size=16, bold=True, color=accent,
        )

        for i in range(1, 4):
            item_top = 1.9 + (i - 1) * 1.5
            item_card = add_rounded_rect(
                slide, f"{prefix}_item_{i}",
                left=left + 0.25, top=item_top, width=5.4, height=1.2,
                fill=WHITE, shadow=False,
            )
            # Icon circle
            add_circle(
                slide, f"{prefix}_icon_{i}",
                left=left + 0.35, top=item_top + 0.3, diameter=0.5,
                fill=accent, text=icon, font_size=12,
            )
            add_textbox(
                slide, f"{prefix}_{i}", f"{{{{{prefix}_{i}}}}}",
                left=left + 1.0, top=item_top + 0.2, width=4.3, height=0.9,
                font_size=12, color=DARK_NAVY,
            )

    add_footer(slide)


def build_slide_09_impact(prs: Presentation) -> None:
    """Slide 9: Impact metrics — navy background, 3 metric cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY_MID)

    add_textbox(
        slide, "impact_title", "{{impact_title}}",
        left=0.5, top=0.3, width=12.33, height=0.6,
        font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, "impact_subtitle", "{{impact_subtitle}}",
        left=0.5, top=0.88, width=12.33, height=0.4,
        font_size=14, color=GRAY_MID, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, "impact_title_accent", 5.9, 0.82, 1.53, 0.04, fill=BRAND_BLUE)

    metric_configs = [
        (1, BRAND_BLUE, 0.5),
        (2, GREEN, 4.72),
        (3, ORANGE, 8.94),
    ]
    card_w = 3.8

    for idx, accent_color, left in metric_configs:
        n = idx
        card = add_rounded_rect(
            slide, f"impact_{n}_card",
            left=left, top=1.5, width=card_w, height=5.2,
            fill=NAVY_LIGHT, shadow=True,
        )
        # Accent top
        add_rounded_rect(
            slide, f"impact_{n}_accent",
            left=left, top=1.5, width=card_w, height=0.15,
            fill=accent_color, shadow=False, adj=10000,
        )
        add_textbox(
            slide, f"impact_{n}_title", f"{{{{impact_{n}_title}}}}",
            left=left + 0.2, top=1.75, width=3.4, height=0.5,
            font_size=13, bold=True, color=WHITE,
        )
        add_textbox(
            slide, f"impact_{n}_metric", f"{{{{impact_{n}_metric}}}}",
            left=left + 0.2, top=2.3, width=3.4, height=0.4,
            font_size=10, color=GRAY_MID,
        )

        # Before value + bar
        add_textbox(
            slide, f"impact_{n}_unit", "ANTES",
            left=left + 0.2, top=2.85, width=1.5, height=0.3,
            font_size=9, color=GRAY_MID,
        )
        add_textbox(
            slide, f"impact_{n}_before_val", f"{{{{impact_{n}_before_val}}}}",
            left=left + 0.2, top=3.2, width=1.8, height=0.55,
            font_size=22, bold=True, color=GRAY_MID,
        )
        before_bar_bg = add_rect(
            slide, f"impact_{n}_before_bar_bg",
            left + 0.2, 3.85, 3.4, 0.22, fill=NAVY_MID,
        )
        add_rect(
            slide, f"impact_{n}_before_bar",
            left + 0.2, 3.85, 1.7, 0.22, fill=GRAY_MID,
        )

        # After value + bar
        add_textbox(
            slide, f"impact_{n}_after_label", "DESPUES",
            left=left + 0.2, top=4.2, width=1.5, height=0.3,
            font_size=9, color=accent_color,
        )
        add_textbox(
            slide, f"impact_{n}_after_val", f"{{{{impact_{n}_after_val}}}}",
            left=left + 0.2, top=4.55, width=1.8, height=0.6,
            font_size=22, bold=True, color=accent_color,
        )
        after_bar_bg = add_rect(
            slide, f"impact_{n}_after_bar_bg",
            left + 0.2, 5.25, 3.4, 0.22, fill=NAVY_MID,
        )
        add_rect(
            slide, f"impact_{n}_after_bar",
            left + 0.2, 5.25, 2.7, 0.22, fill=accent_color,
        )

    # Closing text
    add_textbox(
        slide, "impact_closing", "{{impact_closing}}",
        left=0.5, top=6.85, width=12.33, height=0.4,
        font_size=11, color=GRAY_MID, align=PP_ALIGN.CENTER,
    )


def build_slide_10_comparison(prs: Presentation) -> None:
    """Slide 10: Comparison table."""
    from pptx.util import Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, "comparison_headline", "{{comparison_headline}}",
        left=0.5, top=0.3, width=12.33, height=0.55,
        font_size=22, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, "comparison_accent", 5.7, 0.82, 1.93, 0.05, fill=BRAND_BLUE)

    # Table — 7 rows x 3 cols
    rows, cols = 7, 3
    table_left = Inches(0.5)
    table_top = Inches(1.1)
    table_width = Inches(12.33)
    table_height = Inches(5.8)

    table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table
    slide.shapes[-1].name = "comparison_table"

    col_widths = [Inches(5.5), Inches(3.4), Inches(3.43)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    header_labels = ["Criterio", "Sin BucketsAI", "Con BucketsAI"]
    header_fills = [NAVY_MID, NAVY_MID, BRAND_BLUE]
    for ci, (label, fill_c) in enumerate(zip(header_labels, header_fills)):
        cell = table.cell(0, ci)
        cell.text = label
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_c

    # Data rows
    for ri in range(1, rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PAGE_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

            if ci == 0:
                placeholder = f"{{{{comparison_row_{ri}_criteria}}}}"
                color = DARK_NAVY
                is_bold = True
            elif ci == 1:
                placeholder = f"{{{{comparison_row_{ri}_without}}}}"
                color = TEXT_MUTED
                is_bold = False
            else:
                placeholder = f"{{{{comparison_row_{ri}_with}}}}"
                color = BRAND_BLUE
                is_bold = False

            cell.text = placeholder
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = FONT
            run.font.size = Pt(11)
            run.font.bold = is_bold
            run.font.color.rgb = color

    add_footer(slide)


def build_slide_11_security(prs: Presentation) -> None:
    """Slide 11: Security pillars — 5 cards in a row."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PAGE_BG)

    add_textbox(
        slide, "security_title", "Seguridad y control",
        left=0.5, top=0.3, width=12.33, height=0.55,
        font_size=26, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, "security_subtitle", "{{security_subtitle}}",
        left=0.5, top=0.85, width=12.33, height=0.4,
        font_size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, "security_title_accent", 5.9, 0.82, 1.53, 0.05, fill=BRAND_BLUE)

    card_w = 2.3
    card_gap = 0.21
    start_left = 0.5

    for i in range(1, 6):
        left = start_left + (i - 1) * (card_w + card_gap)
        card = add_rounded_rect(
            slide, f"security_card_{i}",
            left=left, top=1.5, width=card_w, height=5.3,
            fill=WHITE, shadow=True,
        )
        add_rounded_rect(
            slide, f"security_accent_{i}",
            left=left, top=1.5, width=card_w, height=0.12,
            fill=BRAND_BLUE, shadow=False, adj=10000,
        )
        add_circle(
            slide, f"security_num_{i}",
            left=left + (card_w / 2) - 0.3, top=1.75, diameter=0.6,
            fill=BRAND_BLUE, text=str(i), font_size=12,
        )
        add_textbox(
            slide, f"security_label_{i}", f"{{{{security_label_{i}}}}}",
            left=left + 0.1, top=2.5, width=card_w - 0.2, height=0.55,
            font_size=12, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, f"security_desc_{i}", f"{{{{security_desc_{i}}}}}",
            left=left + 0.1, top=3.1, width=card_w - 0.2, height=3.5,
            font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
        )

    add_footer(slide)


def build_slide_12_cta(prs: Presentation) -> None:
    """Slide 12: CTA — dark navy background, contact details."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY_MID)

    # Decorative circles (same style as cover)
    for name, l, t, d in [
        ("cta_deco_1", 11.0, -0.5, 3.0),
        ("cta_deco_2", -0.8, 5.8, 2.5),
    ]:
        c = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d),
        )
        c.name = name
        spPr = c._element.spPr
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", _hex(BRAND_BLUE))
        alpha = etree.SubElement(srgbClr, qn("a:alpha"))
        alpha.set("val", "8000")
        set_no_line(c)

    add_textbox(
        slide, "cta_question", "{{cta_question}}",
        left=0.5, top=1.5, width=12.33, height=1.2,
        font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, "cta_title_accent", 5.5, 2.65, 2.33, 0.05, fill=BRAND_BLUE)

    add_textbox(
        slide, "cta_description", "{{cta_description}}",
        left=1.5, top=2.9, width=10.33, height=0.8,
        font_size=14, color=GRAY_MID, align=PP_ALIGN.CENTER,
    )

    # CTA Button
    btn = add_rounded_rect(
        slide, "cta_button_shape",
        left=5.0, top=3.9, width=3.33, height=0.65,
        fill=BRAND_BLUE, shadow=True, adj=50000,
    )
    add_textbox(
        slide, "cta_button", "{{cta_button}}",
        left=5.0, top=3.93, width=3.33, height=0.6,
        font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )

    # Divider
    add_rect(slide, "cta_divider", 1.5, 4.85, 10.33, 0.03, fill=NAVY_LIGHT)

    # Contact grid
    add_textbox(
        slide, "cta_name", "{{cta_name}}",
        left=0.5, top=5.1, width=4.0, height=0.45,
        font_size=14, bold=True, color=WHITE,
    )
    add_textbox(
        slide, "cta_role", "{{cta_role}}",
        left=0.5, top=5.55, width=4.0, height=0.35,
        font_size=11, color=GRAY_MID,
    )
    add_textbox(
        slide, "cta_email", "{{cta_email}}",
        left=0.5, top=5.95, width=4.5, height=0.35,
        font_size=11, color=BLUE_MEDIUM,
    )
    add_textbox(
        slide, "cta_phone", "{{cta_phone}}",
        left=0.5, top=6.35, width=3.0, height=0.35,
        font_size=11, color=BLUE_MEDIUM,
    )
    add_textbox(
        slide, "cta_url", "buckets-ai.com",
        left=10.5, top=5.1, width=2.5, height=0.35,
        font_size=11, color=GRAY_MID, align=PP_ALIGN.RIGHT,
    )

    # Logo placeholder
    logo_ph = add_rounded_rect(
        slide, "cta_logo",
        left=11.0, top=5.55, width=1.8, height=0.7,
        fill=NAVY_LIGHT, shadow=False,
    )
    add_label_in_shape(logo_ph, "LOGO", font_size=9, color=GRAY_MID, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> None:
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    builders = [
        ("01 Cover",        build_slide_01_cover),
        ("02 Problem",      build_slide_02_problem),
        ("03 Depth",        build_slide_03_depth),
        ("04 Solution",     build_slide_04_solution),
        ("05 How It Works", build_slide_05_how_it_works),
        ("06 Demo",         build_slide_06_demo),
        ("07 Dimensions",   build_slide_07_dimensions),
        ("08 Before/After", build_slide_08_before_after),
        ("09 Impact",       build_slide_09_impact),
        ("10 Comparison",   build_slide_10_comparison),
        ("11 Security",     build_slide_11_security),
        ("12 CTA",          build_slide_12_cta),
    ]

    for label, builder in builders:
        print(f"  Slide {label}")
        builder(prs)

    prs.save(str(OUTPUT_PATH))
    print(f"\nTemplate saved to:\n  {OUTPUT_PATH}")

    # Verify shape names
    print("\nVerifying named shapes per slide:")
    for i, slide in enumerate(prs.slides, start=1):
        names = [s.name for s in slide.shapes if not s.name.startswith("deco_")]
        print(f"  Slide {i:02d}: {len(names)} named shapes")


if __name__ == "__main__":
    main()
