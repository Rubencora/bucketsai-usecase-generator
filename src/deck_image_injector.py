"""
Post-process a Presenton-generated PPTX to inject predefined brand images.

Usage:
    python3 src/deck_image_injector.py input.pptx output.pptx

Injects images from public/deck-assets/ and public/onepager-assets/ into
specific slides based on keyword matching in slide content.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Emu

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)

ASSETS = {
    'logo': os.path.join(PROJECT_DIR, 'public', 'logo.png'),
    'funnel': os.path.join(PROJECT_DIR, 'public', 'deck-assets', 'funnel.png'),
    'frustrated': os.path.join(PROJECT_DIR, 'public', 'deck-assets', 'frustrated-person.png'),
    'comparison': os.path.join(PROJECT_DIR, 'public', 'deck-assets', 'comparison-table.png'),
    'sebastian': os.path.join(PROJECT_DIR, 'public', 'deck-assets', 'sebastian.png'),
    'qr': os.path.join(PROJECT_DIR, 'public', 'deck-assets', 'qr.png'),
    'hero_phones': os.path.join(PROJECT_DIR, 'public', 'onepager-assets', 'hero-phones-clean.png'),
    'chat_phone': os.path.join(PROJECT_DIR, 'public', 'onepager-assets', 'chat-phone-clean.png'),
}

# Slide dimensions (standard 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def get_slide_text(slide):
    """Extract all text from a slide for keyword matching."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text.lower())
    return ' '.join(texts)


def detect_slide_type(slide, index):
    """Detect what type of slide this is based on content keywords."""
    text = get_slide_text(slide)

    # Order matters — check most specific keywords first
    if index == 0:
        return 'cover'

    if any(k in text for k in ['llamado', 'agenda', 'demo personalizada', 'contacto', 'sebastian']):
        return 'cta'

    if any(k in text for k in ['comparaci', 'chatgpt', 'buckets', 'vs chatgpt']):
        if any(k in text for k in ['acceso', 'gobernanza', 'licenciamiento', 'trazabilidad']):
            return 'comparison'

    if any(k in text for k in ['seguridad', 'cumplimiento', 'privacidad', 'aislamiento']):
        return 'security'

    if any(k in text for k in ['impacto', 'métrica', 'velocidad', 'conversión', 'cancelacion']):
        if any(k in text for k in ['antes', 'después', 'cambio', 'antes:']):
            return 'impact'

    if any(k in text for k in ['antes vs', 'antes y después', 'antes']):
        if any(k in text for k in ['después', 'despues']):
            return 'before_after'

    if any(k in text for k in ['dimensi', 'dos dimensiones']):
        return 'dimensions'

    if any(k in text for k in ['ejemplo', 'caso de uso', 'así se', 'mensaje', 'respuesta ia']):
        return 'demo'

    if any(k in text for k in ['cómo funciona', 'como funciona', 'paso', 'step']):
        return 'how_it_works'

    if any(k in text for k in ['solución', 'solucion', 'inteligencia de decisión', 'capacidades']):
        return 'solution'

    if any(k in text for k in ['desafío persiste', 'desafio persiste', 'por qué', 'capacitación no']):
        return 'depth'

    if any(k in text for k in ['problema', 'reto', 'desafío', 'complejidad']):
        return 'problem'

    return 'unknown'


def add_image(slide, image_path, left, top, width=None, height=None):
    """Add an image to a slide at the specified position."""
    if not os.path.exists(image_path):
        print(f'  WARNING: Image not found: {image_path}')
        return None

    kwargs = {'image_file': image_path, 'left': left, 'top': top}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height

    pic = slide.shapes.add_picture(**kwargs)
    return pic


def inject_images(pptx_path, output_path=None):
    """Inject predefined brand images into a Presenton-generated deck."""
    if output_path is None:
        output_path = pptx_path

    prs = Presentation(pptx_path)
    injected = 0

    for i, slide in enumerate(prs.slides):
        slide_type = detect_slide_type(slide, i)
        print(f'  Slide {i+1}: detected as "{slide_type}"')

        if slide_type == 'cover':
            # Logo top-left
            add_image(slide, ASSETS['logo'],
                      left=Inches(0.5), top=Inches(0.4), width=Inches(2.0))
            injected += 1

        elif slide_type == 'problem':
            # Frustrated person — right side
            add_image(slide, ASSETS['frustrated'],
                      left=Inches(8.5), top=Inches(1.5), width=Inches(4.2))
            injected += 1

        elif slide_type == 'depth':
            # Funnel — left side
            add_image(slide, ASSETS['funnel'],
                      left=Inches(0.5), top=Inches(1.5), width=Inches(3.0))
            injected += 1

        elif slide_type == 'solution':
            # Hero phones — right side
            add_image(slide, ASSETS['hero_phones'],
                      left=Inches(8.0), top=Inches(2.0), width=Inches(4.5))
            injected += 1

        elif slide_type == 'demo':
            # Chat phone mockup — right side
            add_image(slide, ASSETS['chat_phone'],
                      left=Inches(8.5), top=Inches(1.0), width=Inches(3.5))
            injected += 1

        elif slide_type == 'comparison':
            # Comparison table image — centered
            add_image(slide, ASSETS['comparison'],
                      left=Inches(1.5), top=Inches(1.8), width=Inches(10.0))
            injected += 1

        elif slide_type == 'cta':
            # Logo top-left + Sebastian photo + QR code
            add_image(slide, ASSETS['logo'],
                      left=Inches(0.5), top=Inches(0.4), width=Inches(2.0))
            add_image(slide, ASSETS['sebastian'],
                      left=Inches(1.0), top=Inches(4.5), width=Inches(1.5))
            add_image(slide, ASSETS['qr'],
                      left=Inches(10.5), top=Inches(4.5), width=Inches(1.8))
            injected += 3

    prs.save(output_path)
    print(f'  Injected {injected} images into {len(prs.slides)} slides')
    print(f'  Saved: {output_path}')
    return output_path


def inject_images_gamma(pptx_path, output_path=None):
    """Gamma mode: only inject logo on cover + CTA info (Sebastian, QR).
    Gamma generates its own AI images, so we skip all other slides."""
    if output_path is None:
        output_path = pptx_path

    prs = Presentation(pptx_path)
    injected = 0

    for i, slide in enumerate(prs.slides):
        slide_type = detect_slide_type(slide, i)
        print(f'  Slide {i+1}: detected as "{slide_type}" (gamma mode)')

        if slide_type == 'cover':
            add_image(slide, ASSETS['logo'],
                      left=Inches(0.5), top=Inches(0.4), width=Inches(2.0))
            injected += 1

        elif slide_type == 'cta':
            add_image(slide, ASSETS['logo'],
                      left=Inches(0.5), top=Inches(0.4), width=Inches(2.0))
            add_image(slide, ASSETS['sebastian'],
                      left=Inches(0.5), top=Inches(5.2), width=Inches(1.3))
            add_image(slide, ASSETS['qr'],
                      left=Inches(11.0), top=Inches(5.2), width=Inches(1.5))
            injected += 3

    prs.save(output_path)
    print(f'  Gamma mode: injected {injected} images into {len(prs.slides)} slides')
    print(f'  Saved: {output_path}')
    return output_path


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: python3 src/deck_image_injector.py input.pptx [output.pptx] [--mode gamma]')
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 and not sys.argv[2].startswith('--') else input_path

    mode = 'full'
    if '--mode' in sys.argv:
        idx = sys.argv.index('--mode')
        if idx + 1 < len(sys.argv):
            mode = sys.argv[idx + 1]

    if mode == 'gamma':
        inject_images_gamma(input_path, output_path)
    else:
        inject_images(input_path, output_path)
